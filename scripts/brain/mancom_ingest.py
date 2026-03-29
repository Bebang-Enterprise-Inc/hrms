"""
BEI Brain: Mancom Report Ingest Pipeline

Downloads reports from Google Chat Mancom space, extracts content,
chunks with semantic awareness, embeds, and stores with full metadata.

Usage:
    python scripts/brain/mancom_ingest.py                     # ingest yesterday's reports
    python scripts/brain/mancom_ingest.py --date 2026-03-24   # ingest specific date
    python scripts/brain/mancom_ingest.py --dry-run            # preview without writing to DB
"""

import os, sys, io, json, hashlib, argparse, time
from datetime import datetime, timedelta, timezone
from pathlib import Path

# Ensure UTF-8 output on Windows
os.environ['PYTHONIOENCODING'] = 'utf-8'
if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8')

from google.oauth2 import service_account
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from openai import OpenAI
from supabase import create_client
from pptx import Presentation
import openpyxl

# ============================================================
# Config
# ============================================================
PROJECT_ROOT = Path(__file__).resolve().parent.parent.parent
SERVICE_ACCOUNT_FILE = PROJECT_ROOT / 'credentials' / 'task-manager-service.json'
MANCOM_SPACE = 'spaces/AAQAVUxMOBM'
USER_ID = '00000000-0000-0000-0000-000000000001'
IMPERSONATE = 'sam@bebang.ph'
EMBED_MODEL = 'text-embedding-3-large'
MAX_CHUNK_CHARS = 1500  # ~375 tokens, safe for embedding

# Secrets from env or Doppler
def get_secret(key):
    val = os.environ.get(key)
    if val:
        return val
    import subprocess
    r = subprocess.run(
        ["C:/Users/Sam/bin/doppler.exe", "secrets", "get", key, "--plain",
         "--project", "bei-erp", "--config", "dev"],
        capture_output=True, text=True
    )
    return r.stdout.strip()


# ============================================================
# Google API helpers
# ============================================================
def get_chat_service():
    creds = service_account.Credentials.from_service_account_file(
        str(SERVICE_ACCOUNT_FILE),
        scopes=['https://www.googleapis.com/auth/chat.messages.readonly']
    ).with_subject(IMPERSONATE)
    return build('chat', 'v1', credentials=creds, cache_discovery=False)


def get_drive_service():
    creds = service_account.Credentials.from_service_account_file(
        str(SERVICE_ACCOUNT_FILE),
        scopes=['https://www.googleapis.com/auth/drive.readonly']
    ).with_subject(IMPERSONATE)
    return build('drive', 'v3', credentials=creds, cache_discovery=False)


def fetch_mancom_messages(chat, date_str):
    """Fetch messages from Mancom space for a given date (PHT)."""
    # Convert PHT date to UTC filter
    utc_start = f"{date_str}T00:00:00+08:00"
    next_day = (datetime.strptime(date_str, '%Y-%m-%d') + timedelta(days=1)).strftime('%Y-%m-%d')
    utc_end = f"{next_day}T00:00:00+08:00"

    messages = []
    page_token = None
    while True:
        kwargs = dict(
            parent=MANCOM_SPACE, pageSize=100,
            filter=f'createTime > "{utc_start}" AND createTime < "{utc_end}"'
        )
        if page_token:
            kwargs['pageToken'] = page_token
        resp = chat.spaces().messages().list(**kwargs).execute()
        messages.extend(resp.get('messages', []))
        page_token = resp.get('nextPageToken')
        if not page_token:
            break
    return messages


def extract_reports(messages):
    """Extract report attachments with Drive file IDs from messages."""
    reports = []
    for msg in messages:
        attachments = msg.get('attachment', [])
        sender = msg.get('sender', {}).get('displayName', 'Unknown')
        text = msg.get('text', '')
        create_time = msg.get('createTime', '')

        for att in attachments:
            drive_ref = att.get('driveDataRef', {})
            file_id = drive_ref.get('driveFileId', '')
            content_name = att.get('contentName', '')

            if file_id and content_name:
                reports.append({
                    'file_id': file_id,
                    'title': content_name,
                    'sender': sender,
                    'message_text': text[:300],
                    'create_time': create_time,
                })
    return reports


def download_file(drive, file_id, out_dir):
    """Download a Google Drive file, returns (local_path, mime_type)."""
    meta = drive.files().get(
        fileId=file_id, fields='name,mimeType',
        supportsAllDrives=True
    ).execute()
    mime = meta['mimeType']
    name = meta.get('name', file_id)

    # Determine export format
    if 'presentation' in mime:
        try:
            request = drive.files().export_media(
                fileId=file_id,
                mimeType='application/vnd.openxmlformats-officedocument.presentationml.presentation'
            )
            ext = '.pptx'
        except Exception:
            request = drive.files().export_media(fileId=file_id, mimeType='text/plain')
            ext = '.txt'
    elif 'spreadsheet' in mime:
        request = drive.files().export_media(
            fileId=file_id,
            mimeType='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
        )
        ext = '.xlsx'
    elif 'document' in mime:
        request = drive.files().export_media(
            fileId=file_id,
            mimeType='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        )
        ext = '.docx'
    else:
        request = drive.files().get_media(fileId=file_id, supportsAllDrives=True)
        ext = '.pdf'

    fh = io.BytesIO()
    downloader = MediaIoBaseDownload(fh, request)
    done = False
    try:
        while not done:
            _, done = downloader.next_chunk()
    except Exception as e:
        if 'too large' in str(e).lower():
            # Fallback to text export for oversized files
            request = drive.files().export_media(fileId=file_id, mimeType='text/plain')
            ext = '.txt'
            fh = io.BytesIO()
            downloader = MediaIoBaseDownload(fh, request)
            done = False
            while not done:
                _, done = downloader.next_chunk()
        else:
            raise

    safe_name = "".join(c if c.isalnum() or c in '-_.' else '_' for c in name)
    out_path = out_dir / f"{safe_name}{ext}"
    fh.seek(0)
    with open(out_path, 'wb') as f:
        f.write(fh.read())

    return out_path, mime


# ============================================================
# Content extraction
# ============================================================
def extract_pptx(filepath):
    """Extract text from PPTX, keeping tables intact."""
    prs = Presentation(str(filepath))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)
            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(cells)
                if rows:
                    # Markdown table
                    parts.append('| ' + ' | '.join(rows[0]) + ' |')
                    parts.append('| ' + ' | '.join(['---'] * len(rows[0])) + ' |')
                    for row in rows[1:]:
                        parts.append('| ' + ' | '.join(row) + ' |')
            if shape.has_chart:
                parts.append('[CHART - visual data]')
        slides.append({
            'slide_num': i,
            'content': '\n'.join(parts) if parts else '(visual-only slide)',
        })
    return slides


def extract_xlsx(filepath):
    """Extract spreadsheet content, keeping tables intact."""
    wb = openpyxl.load_workbook(str(filepath), data_only=True)
    sheets = []
    for name in wb.sheetnames:
        ws = wb[name]
        rows = []
        for row in ws.iter_rows(min_row=1, max_row=min(ws.max_row or 1, 300), values_only=True):
            cells = [str(c) if c is not None else '' for c in row]
            if any(c.strip() for c in cells):
                rows.append(cells)
        if rows:
            sheets.append({'sheet_name': name, 'rows': rows})
    return sheets


def extract_text(filepath):
    """Extract plain text file."""
    with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
        return f.read()


def extract_content(filepath):
    """Route to the right extractor based on file extension."""
    ext = filepath.suffix.lower()
    if ext == '.pptx':
        slides = extract_pptx(filepath)
        parts = []
        for s in slides:
            parts.append(f"## Slide {s['slide_num']}\n\n{s['content']}")
        return '\n\n'.join(parts)
    elif ext == '.xlsx':
        sheets = extract_xlsx(filepath)
        parts = []
        for s in sheets:
            parts.append(f"## Sheet: {s['sheet_name']}\n")
            for row in s['rows'][:200]:
                parts.append(' | '.join(row[:15]))
        return '\n'.join(parts)
    elif ext in ('.txt', '.md'):
        return extract_text(filepath)
    elif ext == '.docx':
        # Basic docx extraction
        try:
            from docx import Document
            doc = Document(str(filepath))
            return '\n\n'.join(p.text for p in doc.paragraphs if p.text.strip())
        except ImportError:
            return extract_text(filepath)
    else:
        return f"(binary file: {filepath.name})"


# ============================================================
# Semantic chunking
# ============================================================
def chunk_content(text, department, date_str, max_chars=MAX_CHUNK_CHARS):
    """
    Chunk text at semantic boundaries (slide headers, sheet headers).
    Keeps tables and data blocks intact where possible.
    """
    # Split on slide/section headers
    import re
    sections = re.split(r'\n(?=## )', text)

    chunks = []
    current = []
    current_len = 0

    for section in sections:
        section_len = len(section)

        # If a single section exceeds max, split it further
        if section_len > max_chars:
            if current:
                chunks.append('\n\n'.join(current))
                current, current_len = [], 0
            # Split large section on double newlines
            paragraphs = section.split('\n\n')
            for para in paragraphs:
                if current_len + len(para) > max_chars and current:
                    chunks.append('\n\n'.join(current))
                    current, current_len = [], 0
                current.append(para)
                current_len += len(para)
            if current:
                chunks.append('\n\n'.join(current))
                current, current_len = [], 0
        elif current_len + section_len > max_chars and current:
            chunks.append('\n\n'.join(current))
            current = [section]
            current_len = section_len
        else:
            current.append(section)
            current_len += section_len

    if current:
        chunks.append('\n\n'.join(current))

    return chunks


# ============================================================
# Detect department from report title/content
# ============================================================
DEPARTMENT_KEYWORDS = {
    'Operations': ['operations', 'ops report', 'ww1', 'ww2', 'gross sales', 'net sales', 'labor cost', 'food cost'],
    'HR': ['hr update', 'human resources', 'government contributions', 'sss', 'hdmf', 'phic', 'last pay'],
    'Marketing': ['marketing', 'creatives', 'content shoot', 'influencer', 'iskrambol', 'poplamig'],
    'Customer Service': ['customer service', 'food panda', 'foodpanda', 'online delivery', 'refund'],
    'R&D': ['r&d', 'status report', 'frozen milk', 'leche flan', 'premix', 'scramble'],
    'Projects': ['project', 'facility', 'design', 'fit-out', 'nso', 'new store opening', 'commissary'],
    'Procurement': ['procurement', 'monitoring', 'savings initiative', 'supplier'],
    'Finance': ['finance', 'budget', 'p&l', 'revenue'],
}


def detect_department(title, content_preview=''):
    """Detect department from title and first ~500 chars of content."""
    text = (title + ' ' + content_preview[:500]).lower()
    scores = {}
    for dept, keywords in DEPARTMENT_KEYWORDS.items():
        score = sum(1 for kw in keywords if kw in text)
        if score > 0:
            scores[dept] = score
    if scores:
        return max(scores, key=scores.get)
    return 'General'


# ============================================================
# Main pipeline
# ============================================================
def ingest_mancom_reports(date_str, dry_run=False):
    print(f"=== BEI Brain Mancom Ingest: {date_str} ===\n")

    # Init services
    openai_client = OpenAI(api_key=get_secret('OPENAI_API_KEY'))
    supabase = create_client(
        get_secret('SUPABASE_URL'),
        get_secret('SUPABASE_SERVICE_ROLE_KEY')
    )
    chat = get_chat_service()
    drive = get_drive_service()

    # Create output directory
    out_dir = PROJECT_ROOT / 'tmp' / f'mancom_{date_str}'
    out_dir.mkdir(parents=True, exist_ok=True)

    # 1. Fetch messages
    print(f"1. Fetching Mancom messages for {date_str}...")
    messages = fetch_mancom_messages(chat, date_str)
    print(f"   Found {len(messages)} messages")

    # 2. Extract report attachments
    reports = extract_reports(messages)
    print(f"   Found {len(reports)} report attachments\n")

    total_chunks = 0
    total_embedded = 0

    for report in reports:
        file_id = report['file_id']
        title = report['title']
        sender = report['sender']

        print(f"--- {title} (by {sender}) ---")

        # 3. Check if already ingested
        existing = supabase.table('documents').select('id').eq('file_id', file_id).execute()
        if existing.data:
            print(f"   SKIP: already ingested (document_id: {existing.data[0]['id']})")
            continue

        # 4. Download
        try:
            filepath, mime = download_file(drive, file_id, out_dir)
            print(f"   Downloaded: {filepath.name} ({filepath.stat().st_size // 1024} KB)")
        except Exception as e:
            print(f"   DOWNLOAD FAILED: {e}")
            continue

        # 5. Extract content
        content = extract_content(filepath)
        if not content or len(content.strip()) < 50:
            print(f"   SKIP: no meaningful content extracted")
            continue

        # 6. Detect department
        department = detect_department(title, content)
        print(f"   Department: {department} | Content: {len(content)} chars")

        # 7. Chunk
        chunks = chunk_content(content, department, date_str)
        print(f"   Chunks: {len(chunks)}")

        if dry_run:
            for i, chunk in enumerate(chunks):
                print(f"   [{i+1}] {len(chunk)} chars: {chunk[:80]}...")
            print()
            continue

        # 8. Create document record
        doc_result = supabase.table('documents').insert({
            'title': title,
            'source_type': 'mancom_report',
            'department': department,
            'author': sender,
            'report_date': date_str,
            'file_id': file_id,
            'chat_space': 'Mancom',
            'chunk_count': len(chunks),
            'metadata': {
                'mime_type': mime,
                'message_text': report['message_text'],
                'create_time': report['create_time'],
                'local_file': str(filepath),
            },
        }).execute()
        document_id = doc_result.data[0]['id']
        print(f"   Document ID: {document_id}")

        # 9. Embed and store chunks
        for i, chunk in enumerate(chunks):
            idem_key = f"mancom:{date_str}:{file_id}:{i}"
            embed_text = (
                f"BEI Mancom Weekly Report\n"
                f"Department: {department}\n"
                f"Date: {date_str}\n"
                f"Author: {sender}\n"
                f"Report: {title}\n\n"
                f"{chunk}"
            )
            if len(embed_text) > 6000:
                embed_text = embed_text[:6000]

            try:
                embedding = openai_client.embeddings.create(
                    input=embed_text, model=EMBED_MODEL
                ).data[0].embedding

                supabase.table('memories').insert({
                    'user_id': USER_ID,
                    'content': chunk[:10000],
                    'embedding': embedding,
                    'embedding_status': 'complete',
                    'source': 'mancom_report',
                    'importance_score': 7,
                    'topic_category': f'mancom_{department.lower().replace(" ", "_")}',
                    'content_hash': hashlib.md5(chunk.encode()).hexdigest(),
                    'idempotency_key': idem_key,
                    'document_id': document_id,
                    'metadata': {
                        'department': department,
                        'author': sender,
                        'date': date_str,
                        'report_title': title,
                        'chunk_index': i,
                        'chunk_total': len(chunks),
                    },
                }).execute()
                total_embedded += 1
            except Exception as e:
                print(f"   chunk {i+1} FAILED: {str(e)[:80]}")

            total_chunks += 1

        print(f"   Embedded: {total_embedded}/{total_chunks}\n")

    print(f"\n=== COMPLETE ===")
    print(f"Reports processed: {len(reports)}")
    print(f"Chunks embedded: {total_embedded}/{total_chunks}")

    return total_embedded


if __name__ == '__main__':
    parser = argparse.ArgumentParser(description='BEI Brain Mancom Ingest')
    parser.add_argument('--date', default=None, help='Date to ingest (YYYY-MM-DD), default: yesterday')
    parser.add_argument('--dry-run', action='store_true', help='Preview without writing to DB')
    args = parser.parse_args()

    if args.date:
        date_str = args.date
    else:
        # Default to yesterday (PHT)
        pht = timezone(timedelta(hours=8))
        date_str = (datetime.now(pht) - timedelta(days=1)).strftime('%Y-%m-%d')

    ingest_mancom_reports(date_str, dry_run=args.dry_run)
