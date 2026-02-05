"""PowerPoint extraction for Knowledge Hub."""

from typing import Dict, Any, List
from pptx import Presentation


def extract_pptx(file_path: str) -> Dict[str, Any]:
    """Extract text content from PowerPoint file.

    Args:
        file_path: Path to the PPTX file.

    Returns:
        Dictionary containing:
        - slides: List of slide content with slide numbers
        - speaker_notes: List of speaker notes with slide numbers
        - metadata: Dictionary with slide_count
    """
    prs = Presentation(file_path)

    result = {
        "slides": [],
        "speaker_notes": [],
        "metadata": {
            "slide_count": len(prs.slides)
        }
    }

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)

        if slide_text:
            result["slides"].append({
                "slide_num": slide_num,
                "content": "\n".join(slide_text)
            })

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes.strip():
                result["speaker_notes"].append({
                    "slide_num": slide_num,
                    "notes": notes.strip()
                })

    return result


def pptx_to_text(file_path: str, include_notes: bool = True) -> str:
    """Convert PPTX to plain text for chunking.

    Args:
        file_path: Path to the PPTX file.
        include_notes: Whether to include speaker notes in output.

    Returns:
        Plain text representation of the presentation.
    """
    extracted = extract_pptx(file_path)

    sections = []
    for slide in extracted["slides"]:
        sections.append(f"--- Slide {slide['slide_num']} ---\n{slide['content']}")

    if include_notes and extracted["speaker_notes"]:
        sections.append("\n--- Speaker Notes ---")
        for note in extracted["speaker_notes"]:
            sections.append(f"Slide {note['slide_num']}: {note['notes']}")

    return "\n\n".join(sections)
